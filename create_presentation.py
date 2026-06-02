from pptx import Presentation
from pptx.util import Inches

prs = Presentation()

slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Monitoramento de Preços — Projeto"
subtitle.text = "IA para identificar tendências de aumento/queda de preços"

slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Resumo"
body = slide.shapes.placeholders[1].text_frame
body.text = "Objetivo: Implementar ETL, treinar modelos e gerar relatórios."
body.add_paragraph().text = "Fontes: API, CSV, Web scraping, PDF"
body.add_paragraph().text = "Entregáveis: CSV/XLSX, SQLite, PDF, Streamlit app"

slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Arquivos importantes"
body = slide.shapes.placeholders[1].text_frame
body.text = "proposta_professor.* (proposta)"
body.add_paragraph().text = "pipeline.py, etl/, app.py, models/, output/"

prs.save('presentation.pptx')
print('Saved presentation.pptx')
